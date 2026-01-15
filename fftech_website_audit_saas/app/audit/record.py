
import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


def export_graphs(audit_id: int, category_scores: dict, out_dir: str) -> str:
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"audit_{audit_id}_categories.png")
    plt.figure(figsize=(8,4))
    plt.bar(category_scores.keys(), category_scores.values(), color='#22c55e')
    plt.title('Category Scores')
    plt.xticks(rotation=30, ha='right')
    plt.tight_layout()
    plt.savefig(path)
    plt.close()
    return path


def export_xlsx(audit_id: int, metrics: dict, category_scores: dict, out_dir: str) -> str:
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"audit_{audit_id}.xlsx")
    with pd.ExcelWriter(path, engine='openpyxl') as writer:
        pd.DataFrame([metrics]).to_excel(writer, index=False, sheet_name='Metrics')
        pd.DataFrame([category_scores]).to_excel(writer, index=False, sheet_name='CategoryScores')
    return path


def export_pptx(audit_id: int, png_chart_path: str, metrics: dict, out_dir: str) -> str:
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"audit_{audit_id}.pptx")
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'FF Tech â€“ Website Audit Summary'
    slide.placeholders[1].text = f'Audit ID: {audit_id}'

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(png_chart_path, Inches(1), Inches(1.2), width=Inches(8))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
    tf.text = 'Key Metrics:'
    for k, v in metrics.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.level = 1

    prs.save(path)
    return path
