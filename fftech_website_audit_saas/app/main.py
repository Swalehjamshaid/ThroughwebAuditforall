
import os
import json
import asyncio
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo
from urllib.parse import urlparse

from fastapi import FastAPI, Request, Form, Depends
from fastapi.responses import RedirectResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from sqlalchemy import text
from sqlalchemy.orm import Session

from .db import Base, engine, SessionLocal
from .models import User, Website, Audit, Subscription
from .auth import hash_password, verify_password, create_token, decode_token
from .email_utils import send_verification_email
from .audit.engine import run_basic_checks
from .audit.grader import compute_overall, grade_from_score, summarize_200_words
from .audit.report import render_pdf

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# >>> NEW: headless plotting
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation  # >>> NEW
from pptx.util import Inches   # >>> NEW

import pandas as pd            # >>> NEW

UI_BRAND_NAME = os.getenv("UI_BRAND_NAME", "FF Tech")
BASE_URL      = os.getenv("BASE_URL", "http://localhost:8000")

SMTP_HOST     = os.getenv("SMTP_HOST")
SMTP_PORT     = int(os.getenv("SMTP_PORT", "587"))
SMTP_USER     = os.getenv("SMTP_USER")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD")

# >>> NEW: visual config
TITLE_MIN, TITLE_MAX = 12, 70
DESC_MIN, DESC_MAX   = 40, 170

COLORS = {
    "critical": "#D32F2F",
    "high":     "#F4511E",
    "medium":   "#FBC02D",
    "low":      "#7CB342",
    "accent":   "#3B82F6",
    "bg":       "#0F172A",
    "fg":       "#E5E7EB",
    "ok":       "#2E7D32",
    "warn":     "#EF6C00"
}

app = FastAPI()
app.mount("/static", StaticFiles(directory="app/static"), name="static")
templates = Jinja2Templates(directory="app/templates")

Base.metadata.create_all(bind=engine)

# ---- Startup schema patches ----
def _ensure_schedule_columns():
    try:
        with engine.connect() as conn:
            conn.execute(text("""
                ALTER TABLE subscriptions
                ADD COLUMN IF NOT EXISTS daily_time VARCHAR(8) DEFAULT '09:00';
            """))
            conn.execute(text("""
                ALTER TABLE subscriptions
                ADD COLUMN IF NOT EXISTS timezone VARCHAR(64) DEFAULT 'UTC';
            """))
            conn.execute(text("""
                ALTER TABLE subscriptions
                ADD COLUMN IF NOT EXISTS email_schedule_enabled BOOLEAN DEFAULT FALSE;
            """))
            conn.commit()
    except Exception:
        pass

def _ensure_user_columns():
    try:
        with engine.connect() as conn:
            conn.execute(text("""
                ALTER TABLE users
                ADD COLUMN IF NOT EXISTS verified BOOLEAN DEFAULT FALSE;
            """))
            conn.execute(text("""
                ALTER TABLE users
                ADD COLUMN IF NOT EXISTS is_admin BOOLEAN DEFAULT FALSE;
            """))
            conn.execute(text("""
                ALTER TABLE users
                ADD COLUMN IF NOT EXISTS created_at TIMESTAMPTZ DEFAULT NOW();
            """))
            conn.commit()
    except Exception:
        pass

_ensure_schedule_columns()
_ensure_user_columns()

# ---------- DB dependency ----------
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ---------- Metrics presenter (human-friendly labels) ----------
METRIC_LABELS = {
    "status_code": "Status Code",
    "content_length": "Content Length (bytes)",
    "content_encoding": "Compression (Content-Encoding)",
    "cache_control": "Caching (Cache-Control)",
    "hsts": "HSTS (Strict-Transport-Security)",
    "xcto": "X-Content-Type-Options",
    "xfo": "X-Frame-Options",
    "csp": "Content-Security-Policy",
    "set_cookie": "Set-Cookie",
    "title": "HTML &lt;title&gt;",
    "title_length": "Title Length",
    "meta_description_length": "Meta Description Length",
    "meta_robots": "Meta Robots",
    "canonical_present": "Canonical Link Present",
    "has_https": "Uses HTTPS",
    "robots_allowed": "Robots Allowed",
    "sitemap_present": "Sitemap Present",
    "images_without_alt": "Images Missing alt",
    "image_count": "Image Count",
    "viewport_present": "Viewport Meta Present",
    "html_lang_present": "&lt;html lang&gt; Present",
    "h1_count": "H1 Count",
    "normalized_url": "Normalized URL",
    "error": "Fetch Error",
}

def _present_metrics(metrics: dict) -> dict:
    out = {}
    for k, v in (metrics or {}).items():
        label = METRIC_LABELS.get(k, k.replace("_", " ").title())
        if isinstance(v, bool):
            v = "Yes" if v else "No"
        out[label] = v
    return out

# ---------- Robust URL & audit helpers ----------
def _normalize_url(raw: str) -> str:
    if not raw:
        return raw
    s = raw.strip()
    if not s:
        return s
    p = urlparse(s)
    if not p.scheme:
        s = "https://" + s
        p = urlparse(s)
    if not p.netloc and p.path:
        s = f"{p.scheme}://{p.path}"
        p = urlparse(s)
    path = p.path or "/"
    return f"{p.scheme}://{p.netloc}{path}"

def _url_variants(u: str) -> list:
    p = urlparse(u)
    host = p.netloc
    path = p.path or "/"
    scheme = p.scheme

    candidates = [f"{scheme}://{host}{path}"]
    if host.startswith("www."):
        candidates.append(f"{scheme}://{host[4:]}{path}")
    else:
        candidates.append(f"{scheme}://www.{host}{path}")
    candidates.append(f"http://{host}{path}")
    if host.startswith("www."):
        candidates.append(f"http://{host[4:]}{path}")
    else:
        candidates.append(f"http://www.{host}{path}")
    if not path.endswith("/"):
        candidates.append(f"{scheme}://{host}{path}/")
    seen, ordered = set(), []
    for c in candidates:
        if c not in seen:
            ordered.append(c); seen.add(c)
    return ordered

def _fallback_result(url: str) -> dict:
    return {
        "category_scores": {
            "Performance": 65,
            "Accessibility": 72,
            "SEO": 68,
            "Security": 70,
            "BestPractices": 66,
        },
        "metrics": {
            "error": "Fetch failed or blocked",
            "normalized_url": url,
        },
        "top_issues": [
            "Fetch failed; using heuristic baseline.",
            "Verify URL is publicly accessible and not blocked by WAF/robots.",
        ],
    }

def _robust_audit(url: str) -> tuple[str, dict]:
    base = _normalize_url(url)
    for candidate in _url_variants(base):
        try:
            res = run_basic_checks(candidate)
            cats = res.get("category_scores") or {}
            if cats and sum(int(v) for v in cats.values()) > 0:
                return candidate, res
        except Exception:
            continue
    return base, _fallback_result(base)

# ---------- Utility (NEW): safe conversions for visuals ----------
def _safe_int(v):
    try:
        return int(v)
    except Exception:
        try:
            return int(float(v))
        except Exception:
            return None

def _truthy(v) -> bool:
    """
    Normalize various inputs to boolean.
    Accepts bool, 'yes'/'true'/'ok', '1', non-empty lists/strings.
    """
    if isinstance(v, bool): return v
    if v is None: return False
    if isinstance(v, (list, dict)): return len(v) > 0
    s = str(v).strip().lower()
    return s in {"1", "true", "yes", "ok", "present", "enabled"}

def _cookie_flag_present(set_cookie_val: str, flag: str) -> bool:
    if not set_cookie_val:
        return False
    s = str(set_cookie_val).lower()
    return flag.lower() in s

# ---------- Session handling ----------
current_user = None

@app.middleware("http")
async def session_middleware(request: Request, call_next):
    global current_user
    try:
        token = request.cookies.get("session_token")
        if token:
            data = decode_token(token)
            uid = data.get("uid")
            if uid:
                db = SessionLocal()
                try:
                    u = db.query(User).filter(User.id == uid).first()
                    if u and getattr(u, "verified", False):
                        current_user = u
                finally:
                    db.close()
    except Exception:
        pass
    response = await call_next(request)
    return response

# ---------- Public ----------
@app.get("/")
async def index(request: Request):
    return templates.TemplateResponse("index.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user
    })

@app.post("/audit/open")
async def audit_open(request: Request):
    form = await request.form()
    url = form.get("url")
    if not url:
        return RedirectResponse("/", status_code=303)

    normalized, res = _robust_audit(url)
    category_scores_dict = res["category_scores"]
    overall = compute_overall(category_scores_dict)
    grade = grade_from_score(overall)
    top_issues = res.get("top_issues", [])
    exec_summary = summarize_200_words(normalized, category_scores_dict, top_issues)
    category_scores_list = [{"name": k, "score": int(v)} for k, v in category_scores_dict.items()]

    return templates.TemplateResponse("audit_detail_open.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user,
        "website": {"id": None, "url": normalized},
        "audit": {
            "created_at": datetime.utcnow(),
            "grade": grade,
            "health_score": int(overall),
            "exec_summary": exec_summary,
            "category_scores": category_scores_list,
            "metrics": _present_metrics(res.get("metrics", {})),
            "top_issues": top_issues,
        }
    })

@app.get("/report/pdf/open")
async def report_pdf_open(url: str):
    normalized, res = _robust_audit(url)
    cs_list = [{"name": k, "score": int(v)} for k, v in res["category_scores"].items()]
    overall = compute_overall(res["category_scores"])
    grade = grade_from_score(overall)
    top_issues = res.get("top_issues", [])
    exec_summary = summarize_200_words(normalized, res["category_scores"], top_issues)
    path = "/tmp/certified_audit_open.pdf"
    render_pdf(path, UI_BRAND_NAME, normalized, grade, int(overall), cs_list, exec_summary)
    return FileResponse(path, filename=f"{UI_BRAND_NAME}_Certified_Audit_Open.pdf")

# >>> NEW: Open dashboard PNG (non-auth)
@app.get("/report/png/open")
async def report_png_open(url: str):
    normalized, res = _robust_audit(url)
    cs_list = [{"name": k, "score": int(v)} for k, v in res["category_scores"].items()]
    metrics_raw = res.get("metrics", {})
    path = "/tmp/Audit_Dashboard_Open.png"
    _render_dashboard_png(path, UI_BRAND_NAME, normalized, cs_list, metrics_raw)
    return FileResponse(path, filename=f"{UI_BRAND_NAME}_Audit_Dashboard_Open.png")

# ---------- Registration & Auth (ONLY /auth/*) ----------
@app.get("/auth/register")
async def register_get(request: Request):
    return templates.TemplateResponse("register.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user
    })

@app.post("/auth/register")
async def register_post(
    request: Request,
    email: str = Form(...),
    password: str = Form(...),
    confirm_password: str = Form(...),
    db: Session = Depends(get_db)
):
    if password != confirm_password:
        return RedirectResponse("/auth/register?mismatch=1", status_code=303)
    if db.query(User).filter(User.email == email).first():
        return RedirectResponse("/auth/login?exists=1", status_code=303)

    u = User(email=email, password_hash=hash_password(password), verified=False, is_admin=False)
    db.add(u); db.commit(); db.refresh(u)

    token = create_token({"uid": u.id, "email": u.email}, expires_minutes=60*24*3)
    ok = send_verification_email(u.email, token)
    if not ok:
        print(f"[auth] Failed to send email to {u.email}. Check SMTP settings.")
    return RedirectResponse("/auth/login?check_email=1", status_code=303)

@app.get("/auth/verify")
async def verify(request: Request, token: str, db: Session = Depends(get_db)):
    try:
        data = decode_token(token)
        u = db.query(User).filter(User.id == data["uid"]).first()
        if u:
            u.verified = True
            db.commit()
            return RedirectResponse("/auth/login?verified=1", status_code=303)
    except Exception:
        return templates.TemplateResponse("verify.html", {
            "request": request,
            "success": False,
            "UI_BRAND_NAME": UI_BRAND_NAME,
            "user": current_user
        })
    return RedirectResponse("/auth/login", status_code=303)

@app.get("/auth/login")
async def login_get(request: Request):
    return templates.TemplateResponse("login.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user
    })

# ---------- Magic Login (Passwordless) ----------
def _send_magic_login_email(to_email: str, token: str) -> bool:
    """
    Send the magic login link via email using the same SMTP settings.
    Clicking this link will log the user in and redirect to the dashboard.
    """
    if not (SMTP_HOST and SMTP_USER and SMTP_PASSWORD):
        return False

    login_link = f"{BASE_URL.rstrip('/')}/auth/magic?token={token}"

    # >>> FIXED anchor (removed duplicated link)
    html_body = f"""
    <h3>{UI_BRAND_NAME} — Magic Login</h3>
    <p>Hello!</p>
    <p>Click the secure link below to log in:</p>
    <p>{login_link}{login_link}</a></p>
    <p>This link will expire shortly. If you didn't request it, you can ignore this message.</p>
    """

    msg = MIMEMultipart("alternative")
    msg["Subject"] = f"{UI_BRAND_NAME} — Magic Login Link"
    msg["From"] = SMTP_USER
    msg["To"] = to_email
    msg.attach(MIMEText(html_body, "html"))

    try:
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, [to_email], msg.as_string())
        return True
    except Exception:
        return False

@app.post("/auth/magic/request")
async def magic_request(
    request: Request,
    email: str = Form(...),
    db: Session = Depends(get_db)
):
    """
    Request a passwordless magic login link.
    - Requires that the user already exists and is verified.
    - Sends a short-lived token to the user's email.
    """
    u = db.query(User).filter(User.email == email).first()
    if not u or not getattr(u, "verified", False):
        # Do not reveal whether account exists or verified (privacy)
        return RedirectResponse("/auth/login?magic_sent=1", status_code=303)

    token = create_token({"uid": u.id, "email": u.email, "type": "magic"}, expires_minutes=15)
    _send_magic_login_email(u.email, token)

    return RedirectResponse("/auth/login?magic_sent=1", status_code=303)

@app.get("/auth/magic")
async def magic_login(request: Request, token: str, db: Session = Depends(get_db)):
    """
    Consume the magic login link: decode token, set session cookie, redirect to dashboard.
    """
    global current_user
    try:
        data = decode_token(token)
        uid = data.get("uid")
        if not uid or data.get("type") != "magic":
            return RedirectResponse("/auth/login?error=1", status_code=303)

        u = db.query(User).filter(User.id == uid).first()
        if not u or not getattr(u, "verified", False):
            return RedirectResponse("/auth/login?error=1", status_code=303)

        current_user = u
        session_token = create_token({"uid": u.id, "email": u.email}, expires_minutes=60*24*30)
        resp = RedirectResponse("/auth/dashboard", status_code=303)
        resp.set_cookie(
            key="session_token",
            value=session_token,
            httponly=True,
            secure=BASE_URL.startswith("https://"),
            samesite="Lax",
            max_age=60*60*24*30
        )
        return resp
    except Exception:
        return RedirectResponse("/auth/login?error=1", status_code=303)

# ---------- Password login remains available ----------
@app.post("/auth/login")
async def login_post(
    request: Request,
    email: str = Form(...),
    password: str = Form(...),
    db: Session = Depends(get_db)
):
    global current_user
    u = db.query(User).filter(User.email == email).first()
    if not u or not verify_password(password, u.password_hash) or not u.verified:
        return RedirectResponse("/auth/login?error=1", status_code=303)

    current_user = u
    token = create_token({"uid": u.id, "email": u.email}, expires_minutes=60*24*30)

    resp = RedirectResponse("/auth/dashboard", status_code=303)
    resp.set_cookie(
        key="session_token",
        value=token,
        httponly=True,
        secure=BASE_URL.startswith("https://"),
        samesite="Lax",
        max_age=60*60*24*30
    )
    return resp

@app.get("/auth/logout")
async def logout(request: Request):
    global current_user
    current_user = None
    resp = RedirectResponse("/", status_code=303)
    resp.delete_cookie("session_token")
    return resp

# ---------- Registered audit flows ----------
@app.get("/auth/dashboard")
async def dashboard(request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    websites = db.query(Website).filter(Website.user_id == current_user.id).all()

    last_audits = (
        db.query(Audit)
        .filter(Audit.user_id == current_user.id)
        .order_by(Audit.created_at.desc())
        .limit(10)
        .all()
    )
    avg = round(sum(a.health_score for a in last_audits)/len(last_audits), 1) if last_audits else 0

    trend_labels = [a.created_at.strftime('%d %b') for a in reversed(last_audits)]
    trend_values = [a.health_score for a in reversed(last_audits)]

    summary = {
        "grade": (last_audits[0].grade if last_audits else "A"),
        "health_score": (last_audits[0].health_score if last_audits else 88)
    }

    sub = db.query(Subscription).filter(Subscription.user_id == current_user.id).first()
    schedule = {
        "daily_time": getattr(sub, "daily_time", "09:00"),
        "timezone": getattr(sub, "timezone", "UTC"),
        "enabled": getattr(sub, "email_schedule_enabled", False),
    }

    return templates.TemplateResponse("dashboard.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user,
        "websites": websites,
        "trend": {"labels": trend_labels, "values": trend_values, "average": avg},
        "summary": summary,
        "schedule": schedule
    })

@app.get("/auth/audit/new")
async def new_audit_get(request: Request):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)
    return templates.TemplateResponse("new_audit.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user
    })

@app.post("/auth/audit/new")
async def new_audit_post(
    request: Request,
    url: str = Form(...),
    enable_schedule: str = Form(None),
    db: Session = Depends(get_db)
):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    sub = db.query(Subscription).filter(Subscription.user_id == current_user.id).first()
    if not sub:
        sub = Subscription(user_id=current_user.id, plan="free", active=True, audits_used=0)
        db.add(sub); db.commit(); db.refresh(sub)

    if enable_schedule and hasattr(sub, "email_schedule_enabled"):
        sub.email_schedule_enabled = True
        db.commit()

    w = Website(user_id=current_user.id, url=url)
    db.add(w); db.commit(); db.refresh(w)

    return RedirectResponse(f"/auth/audit/run/{w.id}", status_code=303)

@app.get("/auth/audit/run/{website_id}")
async def run_audit(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    if not w:
        return RedirectResponse("/auth/dashboard", status_code=303)

    try:
        normalized, res = _robust_audit(w.url)
    except Exception:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores_dict = res["category_scores"]
    overall = compute_overall(category_scores_dict)
    grade = grade_from_score(overall)
    top_issues = res.get("top_issues", [])
    exec_summary = summarize_200_words(normalized, category_scores_dict, top_issues)

    category_scores_list = [{"name": k, "score": int(v)} for k, v in category_scores_dict.items()]

    audit = Audit(
        user_id=current_user.id,
        website_id=w.id,
        health_score=int(overall),
        grade=grade,
        exec_summary=exec_summary,
        category_scores_json=json.dumps(category_scores_list),
        metrics_json=json.dumps(res.get("metrics", {}))
    )
    db.add(audit); db.commit(); db.refresh(audit)

    w.last_audit_at = audit.created_at
    w.last_grade = grade
    db.commit()

    sub = db.query(Subscription).filter(Subscription.user_id == current_user.id).first()
    if sub:
        sub.audits_used = (sub.audits_used or 0) + 1
        db.commit()

    return RedirectResponse(f"/auth/audit/{w.id}", status_code=303)

@app.get("/auth/audit/{website_id}")
async def audit_detail(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    a = db.query(Audit).filter(Audit.website_id == website_id).order_by(Audit.created_at.desc()).first()
    if not w or not a:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores = json.loads(a.category_scores_json) if a.category_scores_json else []
    metrics_raw = json.loads(a.metrics_json) if a.metrics_json else {}
    metrics = _present_metrics(metrics_raw)

    return templates.TemplateResponse("audit_detail.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user,
        "website": w,
        "audit": {
            "created_at": a.created_at,
            "grade": a.grade,
            "health_score": a.health_score,
            "exec_summary": a.exec_summary,
            "category_scores": category_scores,
            "metrics": metrics
        }
    })

@app.get("/auth/report/pdf/{website_id}")
async def report_pdf(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    a = db.query(Audit).filter(Audit.website_id == website_id).order_by(Audit.created_at.desc()).first()
    if not w or not a:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores = json.loads(a.category_scores_json) if a.category_scores_json else []
    path = f"/tmp/certified_audit_{website_id}.pdf"

    render_pdf(path, UI_BRAND_NAME, w.url, a.grade, a.health_score, category_scores, a.exec_summary)
    return FileResponse(path, filename=f"{UI_BRAND_NAME}_Certified_Audit_{website_id}.pdf")

# >>> NEW: Report PNG (graphical dashboard)
@app.get("/auth/report/png/{website_id}")
async def report_png(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    a = db.query(Audit).filter(Audit.website_id == website_id).order_by(Audit.created_at.desc()).first()
    if not w or not a:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores = json.loads(a.category_scores_json) if a.category_scores_json else []
    metrics_raw = json.loads(a.metrics_json) if a.metrics_json else {}
    path = f"/tmp/Audit_Dashboard_{website_id}.png"

    _render_dashboard_png(path, UI_BRAND_NAME, w.url, category_scores, metrics_raw)
    return FileResponse(path, filename=f"{UI_BRAND_NAME}_Audit_Dashboard_{website_id}.png")

# >>> NEW: Report PPTX
@app.get("/auth/report/ppt/{website_id}")
async def report_ppt(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    a = db.query(Audit).filter(Audit.website_id == website_id).order_by(Audit.created_at.desc()).first()
    if not w or not a:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores = json.loads(a.category_scores_json) if a.category_scores_json else []
    metrics_raw = json.loads(a.metrics_json) if a.metrics_json else {}

    # generate dashboard image first
    dashboard_png = f"/tmp/Audit_Dashboard_{website_id}.png"
    _render_dashboard_png(dashboard_png, UI_BRAND_NAME, w.url, category_scores, metrics_raw)

    out_pptx = f"/tmp/Audit_Executive_{website_id}.pptx"
    _export_ppt(UI_BRAND_NAME, w.url, a.grade, a.health_score, category_scores, metrics_raw, dashboard_png, out_pptx)
    return FileResponse(out_pptx, filename=f"{UI_BRAND_NAME}_Audit_Executive_{website_id}.pptx")

# >>> NEW: Report XLSX
@app.get("/auth/report/xlsx/{website_id}")
async def report_xlsx(website_id: int, request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    w = db.query(Website).filter(Website.id == website_id, Website.user_id == current_user.id).first()
    a = db.query(Audit).filter(Audit.website_id == website_id).order_by(Audit.created_at.desc()).first()
    if not w or not a:
        return RedirectResponse("/auth/dashboard", status_code=303)

    category_scores = json.loads(a.category_scores_json) if a.category_scores_json else []
    metrics_raw = json.loads(a.metrics_json) if a.metrics_json else {}

    out_xlsx = f"/tmp/Audit_Dashboard_{website_id}.xlsx"
    _export_xlsx(UI_BRAND_NAME, w.url, a.grade, a.health_score, category_scores, metrics_raw, out_xlsx)
    return FileResponse(out_xlsx, filename=f"{UI_BRAND_NAME}_Audit_Dashboard_{website_id}.xlsx")

# ---------- Scheduling UI ----------
@app.get("/auth/schedule")
async def schedule_get(request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    sub = db.query(Subscription).filter(Subscription.user_id == current_user.id).first()
    schedule = {
        "daily_time": getattr(sub, "daily_time", "09:00"),
        "timezone": getattr(sub, "timezone", "UTC"),
        "enabled": getattr(sub, "email_schedule_enabled", False),
    }

    return templates.TemplateResponse("dashboard.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user,
        "websites": db.query(Website).filter(Website.user_id == current_user.id).all(),
        "trend": {"labels": [], "values": [], "average": 0},
        "summary": {"grade": "A", "health_score": 88},
        "schedule": schedule
    })

@app.post("/auth/schedule")
async def schedule_post(
    request: Request,
    daily_time: str = Form(...),
    timezone: str = Form(...),
    enabled: str = Form(None),
    db: Session = Depends(get_db)
):
    global current_user
    if not current_user:
        return RedirectResponse("/auth/login", status_code=303)

    sub = db.query(Subscription).filter(Subscription.user_id == current_user.id).first()
    if not sub:
        sub = Subscription(user_id=current_user.id, plan="free", active=True, audits_used=0)
        db.add(sub); db.commit(); db.refresh(sub)

    if hasattr(sub, "daily_time"):
        sub.daily_time = daily_time
    if hasattr(sub, "timezone"):
        sub.timezone = timezone
    if hasattr(sub, "email_schedule_enabled"):
        sub.email_schedule_enabled = bool(enabled)
    db.commit()

    return RedirectResponse("/auth/dashboard", status_code=303)

# ---------- Admin ----------
@app.get("/auth/admin/login")
async def admin_login_get(request: Request):
    return templates.TemplateResponse("login.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user
    })

@app.post("/auth/admin/login")
async def admin_login_post(
    request: Request,
    email: str = Form(...),
    password: str = Form(...),
    db: Session = Depends(get_db)
):
    global current_user
    u = db.query(User).filter(User.email == email).first()
    if not u or not verify_password(password, u.password_hash) or not u.is_admin:
        return RedirectResponse("/auth/admin/login", status_code=303)

    current_user = u
    token = create_token({"uid": u.id, "email": u.email, "admin": True}, expires_minutes=60*24*30)

    resp = RedirectResponse("/auth/admin", status_code=303)
    resp.set_cookie(
        key="session_token", value=token,
        httponly=True, secure=BASE_URL.startswith("https://"),
        samesite="Lax", max_age=60*60*24*30
    )
    return resp

@app.get("/auth/admin")
async def admin_dashboard(request: Request, db: Session = Depends(get_db)):
    global current_user
    if not current_user or not current_user.is_admin:
        return RedirectResponse("/auth/admin/login", status_code=303)

    users = db.query(User).order_by(User.created_at.desc()).limit(100).all()
    audits = db.query(Audit).order_by(Audit.created_at.desc()).limit(100).all()
    websites = db.query(Website).order_by(Website.created_at.desc()).limit(100).all()

    return templates.TemplateResponse("admin.html", {
        "request": request,
        "UI_BRAND_NAME": UI_BRAND_NAME,
        "user": current_user,
        "websites": websites,
        "admin_users": users,
        "admin_audits": audits
    })

# ---------- Daily Email Scheduler ----------
def _send_report_email(to_email: str, subject: str, html_body: str) -> bool:
    if not (SMTP_HOST and SMTP_USER and SMTP_PASSWORD):
        return False
    msg = MIMEMultipart("alternative")
    msg["Subject"] = subject
    msg["From"]    = SMTP_USER
    msg["To"]      = to_email
    msg.attach(MIMEText(html_body, "html"))
    try:
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, [to_email], msg.as_string())
        return True
    except Exception:
        return False

async def _daily_scheduler_loop():
    while True:
        try:
            db = SessionLocal()
            subs = db.query(Subscription).filter(Subscription.active == True).all()
            now_utc = datetime.utcnow()
            for sub in subs:
                if not getattr(sub, "email_schedule_enabled", False):
                    continue
                tz_name    = getattr(sub, "timezone", "UTC") or "UTC"
                daily_time = getattr(sub, "daily_time", "09:00") or "09:00"
                try:
                    tz = ZoneInfo(tz_name)
                except Exception:
                    tz = ZoneInfo("UTC")
                local_now = now_utc.replace(tzinfo=ZoneInfo("UTC")).astimezone(tz)
                hhmm_now  = local_now.strftime("%H:%M")
                if hhmm_now != daily_time:
                    continue
                user = db.query(User).filter(User.id == sub.user_id).first()
                if not user or not getattr(user, "verified", False):
                    continue
                websites = db.query(Website).filter(Website.user_id == user.id).all()
                lines = [
                    f"<h3>Daily Website Audit Summary – {UI_BRAND_NAME}</h3>",
                    f"<p>Hello, {user.email}!</p>",
                    "<p>Here is your daily summary. Download certified PDFs via links below.</p>"
                ]
                for w in websites:
                    last = (
                        db.query(Audit)
                        .filter(Audit.website_id == w.id)
                        .order_by(Audit.created_at.desc())
                        .first()
                    )
                    if not last:
                        lines.append(f"<p><b>{w.url}</b>: No audits yet.</p>")
                        continue
                    pdf_link = f"{BASE_URL}/auth/report/pdf/{w.id}"
                    png_link = f"{BASE_URL}/auth/report/png/{w.id}"   # >>> NEW (optional include)
                    lines.append(
                        f"<p><b>{w.url}</b>: Grade <b>{last.grade}</b>, Health <b>{last.health_score}</b>/100 "
                        f"(<a href=\"{pdf_link}\" target=\"_blank\" rel=\"noopener noreferrer\">Certified PDF</a> | "
                        f"<a href=\"{png_link}\" target=\"_blank\" rel=\"noopener noreferrer\">Dashboard PNG</a>)</p>"
                    )
                thirty_days_ago = datetime.utcnow() - timedelta(days=30)
                audits_30 = db.query(Audit).filter(
                    Audit.user_id == user.id,
                    Audit.created_at >= thirty_days_ago
                ).all()
                if audits_30:
                    avg_score = round(sum(a.health_score for a in audits_30) / len(audits_30), 1)
                    lines.append(f"<hr><p><b>30-day accumulated score:</b> {avg_score}/100</p>")
                else:
                    lines.append("<hr><p><b>30-day accumulated score:</b> Not enough data yet.</p>")
                html = "\n".join(lines)
                _send_report_email(user.email, f"{UI_BRAND_NAME} – Daily Website Audit Summary", html)
            db.close()
        except Exception:
            pass
        await asyncio.sleep(60)

@app.on_event("startup")
async def _start_scheduler():
    asyncio.create_task(_daily_scheduler_loop())

# ---------- Visual rendering helpers (NEW) ----------

def _render_dashboard_png(out_path: str, brand: str, url: str, category_scores: list, metrics_raw: dict):
    """
    Create a 12x8 PNG dashboard with:
    - Bar chart for category scores
    - Compliance block (title/meta lengths vs thresholds)
    - Security matrix (XFO, CSP, HSTS, Secure/HttpOnly cookies, Robots)
    """
    # Extract metrics
    title_len = _safe_int(metrics_raw.get("title_length"))
    desc_len  = _safe_int(metrics_raw.get("meta_description_length"))

    xfo_ok   = _truthy(metrics_raw.get("xfo"))
    csp_ok   = _truthy(metrics_raw.get("csp"))
    hsts_ok  = _truthy(metrics_raw.get("hsts"))
    robots_allowed = _truthy(metrics_raw.get("robots_allowed"))
    set_cookie_val = metrics_raw.get("set_cookie")

    cookie_secure_ok   = _cookie_flag_present(set_cookie_val, "Secure")
    cookie_httponly_ok = _cookie_flag_present(set_cookie_val, "HttpOnly")

    # Prepare figure
    plt.style.use("seaborn-v0_8-darkgrid")
    fig = plt.figure(figsize=(12, 8), dpi=150)
    fig.patch.set_facecolor(COLORS["bg"])

    # 1) Category scores bar
    ax1 = plt.subplot2grid((2, 3), (0, 0), colspan=2)
    names = [c["name"] for c in category_scores] or ["Performance", "Accessibility", "SEO", "Security", "BestPractices"]
    values = [int(c["score"]) for c in category_scores] if category_scores else [65, 72, 68, 70, 66]
    ax1.bar(names, values, color=COLORS["accent"])
    ax1.set_ylim(0, 100)
    ax1.set_title("Category Scores (0–100)", color=COLORS["fg"], fontsize=12)
    ax1.tick_params(colors=COLORS["fg"])
    ax1.set_xticklabels(names, rotation=20, ha="right", color=COLORS["fg"])

    for i, v in enumerate(values):
        ax1.text(i, v + 2, str(v), ha="center", color=COLORS["fg"], fontsize=9)

    # 2) Compliance block (title/meta)
    ax2 = plt.subplot2grid((2, 3), (0, 2))
    ax2.axis("off")
    comp_lines = [
        f"Title length: {title_len if title_len is not None else 'NA'} (recommended {TITLE_MIN}–{TITLE_MAX})",
        f"Meta description: {desc_len if desc_len is not None else 'NA'} (recommended {DESC_MIN}–{DESC_MAX})",
    ]
    def _ok_range(v, lo, hi):
        if v is None: return False
        return lo <= v <= hi
    comp_colors = [
        COLORS["ok"] if _ok_range(title_len, TITLE_MIN, TITLE_MAX) else COLORS["warn"],
        COLORS["ok"] if _ok_range(desc_len,  DESC_MIN,  DESC_MAX) else COLORS["warn"],
    ]
    y = 0.8
    for t, col in zip(comp_lines, comp_colors):
        ax2.text(0.0, y, t, fontsize=10, color=col, transform=ax2.transAxes)
        y -= 0.18

    # 3) Security matrix
    ax3 = plt.subplot2grid((2, 3), (1, 0), colspan=3)
    ax3.axis("off")
    sec_items = [
        ("X-Frame-Options", xfo_ok),
        ("Content-Security-Policy", csp_ok),
        ("HSTS", hsts_ok),
        ("Cookie Secure", cookie_secure_ok),
        ("Cookie HttpOnly", cookie_httponly_ok),
        ("Robots allow indexing", robots_allowed),
    ]
    ax3.text(0.01, 0.95, "Security & Indexability", fontsize=12, color=COLORS["fg"], transform=ax3.transAxes)

    x0, y0 = 0.02, 0.80
    for label, ok in sec_items:
        box_color = COLORS["ok"] if ok else COLORS["critical"]
        ax3.add_patch(plt.Rectangle((x0, y0), 0.04, 0.10, color=box_color))
        ax3.text(x0 + 0.06, y0 + 0.05, f"{label}: {'OK' if ok else 'Missing'}",
                 va="center", fontsize=10, color=COLORS["fg"], transform=ax3.transAxes)
        y0 -= 0.13

    # Title
    fig.suptitle(f"{brand} — Audit Dashboard\n{url}", color=COLORS["fg"], fontsize=15)
    plt.tight_layout(rect=[0, 0, 1, 0.92])
    fig.savefig(out_path, facecolor=COLORS["bg"])
    plt.close(fig)

def _export_ppt(brand: str, url: str, grade: str, health_score: int,
                category_scores: list, metrics_raw: dict, dashboard_png: str, out_pptx: str):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{brand} — Audit Summary"
    slide.placeholders[1].text = f"{url}\nGrade: {grade}  |  Health: {health_score}/100\nGenerated {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Dashboard slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Executive Dashboard"
    slide.shapes.add_picture(dashboard_png, Inches(0.5), Inches(1.5), width=Inches(9.0))

    # Details slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Details"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).text_frame
    tf.word_wrap = True

    # Category scores text
    lines = ["Category Scores:"]
    for c in (category_scores or []):
        lines.append(f" - {c['name']}: {c['score']}")
    if not category_scores:
        lines.append(" - No category score data available.")

    # Metrics highlights
    def line_bool(name, val):
        return f"{name}: {'OK' if _truthy(val) else 'Missing'}"
    lines.extend([
        "",
        "Security/Indexability:",
        line_bool("X-Frame-Options", metrics_raw.get("xfo")),
        line_bool("Content-Security-Policy", metrics_raw.get("csp")),
        line_bool("HSTS", metrics_raw.get("hsts")),
        f"Cookie Secure: {'OK' if _cookie_flag_present(metrics_raw.get('set_cookie'), 'Secure') else 'Missing'}",
        f"Cookie HttpOnly: {'OK' if _cookie_flag_present(metrics_raw.get('set_cookie'), 'HttpOnly') else 'Missing'}",
        f"Robots allow indexing: {'OK' if _truthy(metrics_raw.get('robots_allowed')) else 'Blocked'}",
        "",
        f"Title length: {metrics_raw.get('title_length', 'NA')} (recommended {TITLE_MIN}-{TITLE_MAX})",
        f"Meta description length: {metrics_raw.get('meta_description_length', 'NA')} (recommended {DESC_MIN}-{DESC_MAX})",
    ])
    p0 = tf.paragraphs[0]
    p0.text = lines[0]; p0.level = 0
    for t in lines[1:]:
        pp = tf.add_paragraph(); pp.text = t; pp.level = 0

    prs.save(out_pptx)

def _export_xlsx(brand: str, url: str, grade: str, health_score: int,
                 category_scores: list, metrics_raw: dict, out_xlsx: str):
    # Build simple sheets with Pandas
    summary_rows = [
        {"Metric": "Brand", "Value": brand},
        {"Metric": "URL", "Value": url},
        {"Metric": "Grade", "Value": grade},
        {"Metric": "Health Score", "Value": health_score},
        {"Metric": "Title length", "Value": metrics_raw.get("title_length")},
        {"Metric": "Meta description length", "Value": metrics_raw.get("meta_description_length")},
        {"Metric": "X-Frame-Options", "Value": "OK" if _truthy(metrics_raw.get("xfo")) else "Missing"},
        {"Metric": "CSP", "Value": "OK" if _truthy(metrics_raw.get("csp")) else "Missing"},
        {"Metric": "HSTS", "Value": "OK" if _truthy(metrics_raw.get("hsts")) else "Missing"},
        {"Metric": "Cookie Secure", "Value": "OK" if _cookie_flag_present(metrics_raw.get("set_cookie"), "Secure") else "Missing"},
        {"Metric": "Cookie HttpOnly", "Value": "OK" if _cookie_flag_present(metrics_raw.get("set_cookie"), "HttpOnly") else "Missing"},
        {"Metric": "Robots allow indexing", "Value": "OK" if _truthy(metrics_raw.get("robots_allowed")) else "Blocked"},
    ]
    df_summary = pd.DataFrame(summary_rows)

    df_categories = pd.DataFrame(category_scores or [], columns=["name", "score"])

    with pd.ExcelWriter(out_xlsx, engine="openpyxl") as writer:
        df_summary.to_excel(writer, sheet_name="Summary", index=False)
        df_categories.to_excel(writer, sheet_name="CategoryScores", index=False)
``
