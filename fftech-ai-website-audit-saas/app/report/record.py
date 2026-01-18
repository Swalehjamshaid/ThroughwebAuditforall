
from typing import Dict, Any
import io
import pandas as pd
from matplotlib import pyplot as plt
from pptx import Presentation
from pptx.util import Inches


def export_png(metrics: Dict[str, Any]) -> bytes:
    cats = {}
    for m in metrics.values():
        cat = m.get('category','Other')
        s = m.get('score')
        if s is None: continue
        cats.setdefault(cat, []).append(s)
    avg = {k: sum(v)/len(v) for k,v in cats.items() if v}
    fig, ax = plt.subplots(figsize=(8,4))
    ax.bar(list(avg.keys()), list(avg.values()), color='#0b5ed7')
    ax.set_title('Category Scores')
    ax.set_ylabel('Score')
    plt.xticks(rotation=45, ha='right')
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png')
    plt.close(fig)
    return buf.getvalue()


def export_xlsx(audit: Dict[str, Any]) -> bytes:
    df = pd.DataFrame([
        {'id': int(k), 'name': v.get('name'), 'category': v.get('category'), 'score': v.get('score'), 'value': v.get('value')}
        for k,v in audit['metrics'].items()
    ])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Metrics')
    return buf.getvalue()


def export_pptx(audit: Dict[str, Any]) -> bytes:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "FF Tech – AI Website Audit"
    slide.placeholders[1].text = f"Score: {audit['overall']['score']} • Grade: {audit['overall']['grade']} • Coverage: {audit['overall']['coverage']}%"

    img_bytes = export_png(audit['metrics'])
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(1); top = Inches(1)
    slide2.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=Inches(8))

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
