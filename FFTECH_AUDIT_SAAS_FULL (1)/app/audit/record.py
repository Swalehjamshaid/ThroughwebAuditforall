import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

plt.switch_backend('Agg')

def save_charts_png(out_dir: str, metrics: dict, scores: dict):
    os.makedirs(out_dir, exist_ok=True)
    plt.figure(figsize=(6,3))
    names=list(scores.keys()); vals=list(scores.values())
    plt.barh(names, vals, color='#2e86de'); plt.xlim(0,100); plt.tight_layout()
    fp=os.path.join(out_dir,'scores.png'); plt.savefig(fp, dpi=160); plt.close(); return {'scores': fp}


def save_excel(path: str, metrics: dict, scores: dict):
    with pd.ExcelWriter(path, engine='openpyxl') as xw:
        pd.DataFrame([metrics]).to_excel(xw, sheet_name='metrics', index=False)
        pd.DataFrame([scores]).to_excel(xw, sheet_name='scores', index=False)
    return path


def save_pptx(path: str, charts: dict, pdf_path: str | None = None):
    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = 'FF Tech Website Audit Summary'; s.placeholders[1].text = 'Auto-generated deck'
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); s2.shapes.add_picture(charts['scores'], Inches(1), Inches(1), height=Inches(4))
    prs.save(path); return path
